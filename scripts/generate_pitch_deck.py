"""Generates the Sentinel AI hackathon pitch deck as a .pptx file.

Not shipped library code -- a one-off content generator for the presentation. Run:

    pip install python-pptx
    python scripts/generate_pitch_deck.py

Produces sentinel-ai/Sentinel_AI_Pitch.pptx in the repo root.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

# -- palette --------------------------------------------------------------
BG_DARK = RGBColor(0x0B, 0x11, 0x20)
CARD = RGBColor(0x15, 0x1D, 0x30)
ACCENT_GREEN = RGBColor(0x34, 0xD3, 0x99)
ACCENT_AMBER = RGBColor(0xF5, 0xA6, 0x23)
ACCENT_RED = RGBColor(0xEF, 0x44, 0x44)
TEXT_WHITE = RGBColor(0xF1, 0xF5, 0xF9)
TEXT_MUTED = RGBColor(0x94, 0xA3, 0xB8)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _bg(slide, color=BG_DARK) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _textbox(slide, left, top, width, height, text, *, size=18, color=TEXT_WHITE,
             bold=False, italic=False, align=PP_ALIGN.LEFT, font="Calibri", anchor=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    if anchor is not None:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return box


def _bullets(slide, left, top, width, height, items, *, size=16, color=TEXT_WHITE,
             gap=6, bold_lead=True):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        if isinstance(item, tuple):
            lead, rest = item
            r1 = p.add_run()
            r1.text = f"{lead}  "
            r1.font.bold = bold_lead
            r1.font.size = Pt(size)
            r1.font.color.rgb = ACCENT_GREEN
            r2 = p.add_run()
            r2.text = rest
            r2.font.size = Pt(size)
            r2.font.color.rgb = color
        else:
            r = p.add_run()
            r.text = f"•  {item}"
            r.font.size = Pt(size)
            r.font.color.rgb = color
    return box


def _card(slide, left, top, width, height, color=CARD):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # 1 = rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    shape.shadow.inherit = False
    return shape


def _kicker_title(slide, kicker, title, *, title_size=32):
    _textbox(slide, Inches(0.6), Inches(0.35), Inches(8), Inches(0.4), kicker.upper(),
             size=14, color=ACCENT_GREEN, bold=True)
    _textbox(slide, Inches(0.6), Inches(0.68), Inches(11.5), Inches(0.9), title,
             size=title_size, color=TEXT_WHITE, bold=True)


def build() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # ---------------- Slide 1: Title ----------------
    s = prs.slides.add_slide(blank)
    _bg(s)
    _textbox(s, Inches(0.8), Inches(2.5), Inches(11.5), Inches(0.5), "AUTONOMOUS TRAFFIC INTELLIGENCE",
              size=16, color=ACCENT_GREEN, bold=True)
    _textbox(s, Inches(0.8), Inches(2.9), Inches(11.5), Inches(1.4), "Sentinel AI",
              size=60, color=TEXT_WHITE, bold=True)
    _textbox(s, Inches(0.8), Inches(4.1), Inches(10.5), Inches(0.9),
              "A safety-provable, explainable multi-agent system that watches an intersection\n"
              "and drives its signals in real time.",
              size=20, color=TEXT_MUTED)
    _textbox(s, Inches(0.8), Inches(6.7), Inches(8), Inches(0.4),
              "11 agents  ·  209 tests  ·  0 safety violations  ·  72% less waiting",
              size=14, color=ACCENT_AMBER, bold=True)

    # ---------------- Slide 2: The Problem ----------------
    s = prs.slides.add_slide(blank)
    _bg(s)
    _kicker_title(s, "The Problem", "Traffic lights are blind")
    _bullets(s, Inches(0.7), Inches(1.9), Inches(6.0), Inches(4.5), [
        ("Fixed timers", "run on a clock, not on traffic — green burns on empty roads while the other side backs up."),
        ("Wasted everything", "fuel, time, emissions — congestion costs cities billions every year."),
        ("No visibility", "operators can't ask a traffic light *why* it did what it did."),
        ("No safety story", "most \"smart\" signal pilots can't prove they won't cause a conflict."),
    ], size=18, gap=18)
    card = _card(s, Inches(7.3), Inches(1.9), Inches(5.3), Inches(4.5))
    _textbox(s, Inches(7.6), Inches(2.1), Inches(4.7), Inches(0.5), "What a fixed timer sees",
              size=14, color=TEXT_MUTED, bold=True)
    _textbox(s, Inches(7.6), Inches(2.7), Inches(4.7), Inches(3.4),
              "North-South: GREEN for 30s\n(regardless of demand)\n\n"
              "East-West: RED for 30s\n(even if empty)\n\n"
              "→ no sensing. no reasoning. no adaptation.",
              size=16, color=ACCENT_RED)

    # ---------------- Slide 3: The Solution ----------------
    s = prs.slides.add_slide(blank)
    _bg(s)
    _kicker_title(s, "The Solution", "An AI that sees, reasons, and explains")
    _bullets(s, Inches(0.7), Inches(1.9), Inches(11.5), Inches(4.6), [
        ("Sees", "detects and tracks every vehicle per approach, estimates queue length, wait time, density."),
        ("Reasons", "a multi-objective policy weighs queue, wait, fairness, pedestrians, and predicted congestion."),
        ("Acts safely", "a provably-safe controller enforces min/max green and full clearance — no exceptions."),
        ("Explains", "an LLM narrates every decision in plain English, grounded only in real computed facts."),
        ("Proves it", "benchmarked against a fixed timer on identical traffic: 72% less average wait."),
    ], size=19, gap=16)

    # ---------------- Slide 4: Architecture ----------------
    s = prs.slides.add_slide(blank)
    _bg(s)
    _kicker_title(s, "Architecture", "Eleven agents, one closed loop", title_size=30)
    planes = [
        ("PERCEPTION PLANE", ACCENT_GREEN, ["Vision", "Tracking", "Movement Analysis"]),
        ("COGNITION PLANE", ACCENT_AMBER, ["Traffic Memory", "Prediction", "Decision", "Incident Detection", "Explainability"]),
        ("CONTROL & OPS", RGBColor(0x60, 0xA5, 0xFA), ["Signal Controller", "Orchestrator", "Dashboard"]),
    ]
    x = Inches(0.6)
    w = Inches(4.0)
    for label, color, items in planes:
        _card(s, x, Inches(1.9), w, Inches(4.6))
        bar = slide_bar = s.shapes.add_shape(1, x, Inches(1.9), w, Inches(0.5))
        bar.fill.solid(); bar.fill.fore_color.rgb = color; bar.line.color.rgb = color
        tf = bar.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = label; r.font.bold = True; r.font.size = Pt(14)
        r.font.color.rgb = BG_DARK
        _bullets(s, x + Inches(0.25), Inches(2.6), w - Inches(0.5), Inches(3.7),
                 items, size=16, gap=14)
        x += w + Inches(0.35)
    _textbox(s, Inches(0.6), Inches(6.75), Inches(12), Inches(0.5),
              "state.updated → prediction.updated → decision.made → signal.changed → explanation.generated",
              size=13, color=TEXT_MUTED, italic=True)

    # ---------------- Slide 5: Safety story ----------------
    s = prs.slides.add_slide(blank)
    _bg(s)
    _kicker_title(s, "The Safety Story", "The LLM never touches the light")
    _card(s, Inches(0.7), Inches(1.9), Inches(5.6), Inches(4.6))
    _textbox(s, Inches(1.0), Inches(2.1), Inches(5.0), Inches(0.5), "Tier 1 — Decision Policy",
              size=18, color=ACCENT_AMBER, bold=True)
    _bullets(s, Inches(1.0), Inches(2.7), Inches(5.0), Inches(3.6), [
        "Multi-objective utility score per axis",
        "Can be wrong. Can be replaced. Can even be an RL model later.",
        "Produces intent only — never touches the actuator",
    ], size=16, gap=14)
    _card(s, Inches(6.6), Inches(1.9), Inches(6.0), Inches(4.6))
    _textbox(s, Inches(6.9), Inches(2.1), Inches(5.4), Inches(0.5), "Tier 2 — Signal Controller",
              size=18, color=ACCENT_GREEN, bold=True)
    _bullets(s, Inches(6.9), Inches(2.7), Inches(5.4), Inches(3.6), [
        "Enforces min/max green, mandatory yellow + all-red clearance",
        "Rejects any illegal phase transition, regardless of policy input",
        "sentinel_safety_violations_total metric — verified at 0 across 200+ tests",
    ], size=16, gap=14)
    _textbox(s, Inches(0.7), Inches(6.75), Inches(11.5), Inches(0.5),
              "Explainability Agent narrates the decision afterward — out of the loop, fails safe to a template.",
              size=14, color=TEXT_MUTED, italic=True)

    # ---------------- Slide 6: Proof / Results ----------------
    s = prs.slides.add_slide(blank)
    _bg(s)
    _kicker_title(s, "Proof, Not Just a Pitch", "Measured against a real baseline")
    metrics = [
        ("72%", "less average wait", ACCENT_GREEN),
        ("11/11", "agents implemented & running", ACCENT_AMBER),
        ("209", "automated tests passing", RGBColor(0x60, 0xA5, 0xFA)),
        ("0", "safety violations, ever", ACCENT_RED),
    ]
    x = Inches(0.6)
    w = Inches(2.95)
    for value, label, color in metrics:
        _card(s, x, Inches(2.0), w, Inches(2.2))
        _textbox(s, x, Inches(2.25), w, Inches(1.0), value, size=44, color=color, bold=True, align=PP_ALIGN.CENTER)
        _textbox(s, x + Inches(0.15), Inches(3.25), w - Inches(0.3), Inches(0.8), label,
                  size=14, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
        x += w + Inches(0.2)
    _textbox(s, Inches(0.6), Inches(4.7), Inches(11.8), Inches(0.4),
              "Fixed timer vs. Sentinel AI — identical seeded traffic, same demand, same duration:",
              size=15, color=TEXT_MUTED)
    _card(s, Inches(0.6), Inches(5.2), Inches(5.6), Inches(1.7))
    _textbox(s, Inches(0.9), Inches(5.4), Inches(5), Inches(0.4), "Fixed timer", size=14, color=TEXT_MUTED)
    _textbox(s, Inches(0.9), Inches(5.7), Inches(5), Inches(0.9), "50.2s avg wait  ·  29 max queue",
              size=20, color=ACCENT_RED, bold=True)
    _card(s, Inches(6.4), Inches(5.2), Inches(5.6), Inches(1.7))
    _textbox(s, Inches(6.7), Inches(5.4), Inches(5), Inches(0.4), "Sentinel AI", size=14, color=TEXT_MUTED)
    _textbox(s, Inches(6.7), Inches(5.7), Inches(5), Inches(0.9), "13.8s avg wait  ·  11 max queue",
              size=20, color=ACCENT_GREEN, bold=True)

    # ---------------- Slide 7: Tech stack ----------------
    s = prs.slides.add_slide(blank)
    _bg(s)
    _kicker_title(s, "Under the Hood", "Real engineering, not a demo hack", title_size=30)
    cols = [
        ("Backend", ["Python 3.11", "FastAPI + Uvicorn + WebSocket", "asyncio event-driven agents", "Pydantic v2 everywhere"]),
        ("AI / Perception", ["YOLO (Ultralytics)", "Custom IoU multi-object tracker", "Claude LLM narrator (out of loop)", "Least-squares trend forecasting"]),
        ("Infra & Ops", ["RabbitMQ / Redis Streams", "PostgreSQL", "Docker + docker-compose", "Prometheus + Grafana", "GitHub Actions CI"]),
    ]
    x = Inches(0.6)
    w = Inches(4.0)
    for title, items in cols:
        _card(s, x, Inches(1.9), w, Inches(4.7))
        _textbox(s, x + Inches(0.25), Inches(2.05), w - Inches(0.5), Inches(0.5), title,
                  size=17, color=ACCENT_GREEN, bold=True)
        _bullets(s, x + Inches(0.25), Inches(2.65), w - Inches(0.5), Inches(3.7), items, size=15, gap=12)
        x += w + Inches(0.35)

    # ---------------- Slide 8: Deployment ----------------
    s = prs.slides.add_slide(blank)
    _bg(s)
    _kicker_title(s, "Ready to Ship", "Hybrid deployment, day one")
    _card(s, Inches(0.7), Inches(2.0), Inches(5.6), Inches(4.2))
    _textbox(s, Inches(1.0), Inches(2.2), Inches(5.0), Inches(0.5), "Vercel", size=20, color=ACCENT_GREEN, bold=True)
    _bullets(s, Inches(1.0), Inches(2.8), Inches(5.0), Inches(3.2), [
        "Static live dashboard",
        "Thin API proxy via vercel.json",
        "Zero build step",
    ], size=16, gap=14)
    _card(s, Inches(6.6), Inches(2.0), Inches(6.0), Inches(4.2))
    _textbox(s, Inches(6.9), Inches(2.2), Inches(5.4), Inches(0.5), "GPU Host", size=20, color=ACCENT_AMBER, bold=True)
    _bullets(s, Inches(6.9), Inches(2.8), Inches(5.4), Inches(3.2), [
        "Full agent fleet + perception",
        "SUMO microsimulation bridge",
        "RabbitMQ / Redis / Postgres",
        "Dockerized, CI-built on every push",
    ], size=16, gap=14)

    # ---------------- Slide 9: Close ----------------
    s = prs.slides.add_slide(blank)
    _bg(s)
    _textbox(s, Inches(0.8), Inches(2.6), Inches(11.5), Inches(1.2),
              "This isn't a concept.", size=40, color=TEXT_WHITE, bold=True)
    _textbox(s, Inches(0.8), Inches(3.6), Inches(11.5), Inches(1.0),
              "It's running. It's safe. It's ready for a real city.",
              size=24, color=ACCENT_GREEN, bold=True)
    _textbox(s, Inches(0.8), Inches(5.2), Inches(10), Inches(0.5), "Sentinel AI — Autonomous Traffic Intelligence",
              size=14, color=TEXT_MUTED)

    return prs


def main() -> None:
    prs = build()
    out = Path(__file__).resolve().parents[1] / "Sentinel_AI_Pitch.pptx"
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    main()
